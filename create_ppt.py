#!/usr/bin/env python3
"""
Battugs Trading System - Combined Entry Setup PPT
Degi + Fractal + TBM нэгдсэн арга барил
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
BG_DARK = RGBColor(0x1a, 0x1a, 0x2e)
BG_CARD = RGBColor(0x16, 0x21, 0x3e)
BG_CARD2 = RGBColor(0x1f, 0x2b, 0x4d)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xAA)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
GREEN = RGBColor(0x08, 0x99, 0x81)
RED = RGBColor(0xF2, 0x36, 0x45)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
BLUE = RGBColor(0x54, 0x9B, 0xFF)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
PURPLE = RGBColor(0xBB, 0x86, 0xFC)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf

def add_para(tf, text, font_size=16, color=WHITE, bold=False, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = space_before
    return p

def add_arrow(slide, left, top, width, height, color=GREEN):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ══════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

# Title decoration line
add_shape(slide, Inches(3), Inches(2.5), Inches(10), Pt(3), GOLD)

add_text(slide, Inches(1), Inches(2.8), Inches(14), Inches(1.5),
         "COMBINED ENTRY SYSTEM", 48, GOLD, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.2), Inches(14), Inches(0.8),
         "Degi  +  Fractal  +  TBM", 32, WHITE, False, PP_ALIGN.CENTER)

add_shape(slide, Inches(3), Inches(5.2), Inches(10), Pt(3), GOLD)

add_text(slide, Inches(1), Inches(5.8), Inches(14), Inches(0.6),
         "3 арга барилыг нэгтгэсэн нарийн entry setup", 20, GRAY, False, PP_ALIGN.CENTER)

# 3 boxes at bottom
for i, (label, col) in enumerate([("DEGI", CYAN), ("FRACTAL", GREEN), ("TBM", ORANGE)]):
    x = Inches(3.5 + i * 3.5)
    card = add_shape(slide, x, Inches(6.8), Inches(2.5), Inches(1.2), BG_CARD, col)
    add_text(slide, x, Inches(7.05), Inches(2.5), Inches(0.5), label, 22, col, True, PP_ALIGN.CENTER)
    subtexts = ["Яагаад?", "Хаана?", "Хэзээ?"]
    add_text(slide, x, Inches(7.45), Inches(2.5), Inches(0.4), subtexts[i], 16, GRAY, False, PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(8.3), Inches(14), Inches(0.5),
         "by Battugs  |  2026", 14, RGBColor(0x66, 0x66, 0x66), False, PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 2: Overview - 3 Layer System
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(14), Inches(0.8),
         "СИСТЕМИЙН БҮТЭЦ — 3 ДАВХАР ШҮҮЛТҮҮР", 36, GOLD, True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(6), Pt(2), GOLD)

# Layer 1 - DEGI
y = 1.6
card = add_shape(slide, Inches(0.8), Inches(y), Inches(14.4), Inches(2), BG_CARD, CYAN)
add_text(slide, Inches(1.2), Inches(y + 0.15), Inches(3), Inches(0.5),
         "ДАВХАР 1: DEGI — ЯАГААД?", 22, CYAN, True)
tf = add_text(slide, Inches(1.2), Inches(y + 0.65), Inches(6.5), Inches(1.2),
              "HTF (1H, 4H) дээр чиглэл тодорхойлох", 16, LIGHT_GRAY)
add_para(tf, "Liquidity sweep + IDM хамалт хүлээх", 16, LIGHT_GRAY)
add_para(tf, "Choch / BOS баталгаажуулах", 16, LIGHT_GRAY)
add_para(tf, "Entry Grade: A.1 (хамгийн сайн) → A → Potential A.1", 16, WHITE, True)

# Right side summary
tf2 = add_text(slide, Inches(8.5), Inches(y + 0.5), Inches(6), Inches(1.3),
               "Гол зорилго:", 16, CYAN, True)
add_para(tf2, '"Том цаг юу хэлж байна?"', 18, WHITE, True)
add_para(tf2, "Трэндийн чиглэл + liquidity байршил", 15, GRAY)
add_para(tf2, "= Арилжааны ЧИГ тодорхойлно", 15, GRAY)

# Arrow
add_arrow(slide, Inches(7.7), Inches(y + 2.05), Inches(0.5), Inches(0.4), CYAN)

# Layer 2 - FRACTAL
y = 4.2
card = add_shape(slide, Inches(0.8), Inches(y), Inches(14.4), Inches(2), BG_CARD, GREEN)
add_text(slide, Inches(1.2), Inches(y + 0.15), Inches(5), Inches(0.5),
         "ДАВХАР 2: FRACTAL — ХААНА?", 22, GREEN, True)
tf = add_text(slide, Inches(1.2), Inches(y + 0.65), Inches(6.5), Inches(1.2),
              "Бүс тодорхойлох (ranging zone)", 16, LIGHT_GRAY)
add_para(tf, "Шилжүүлэгч хөдөлгөөн таних (displacement)", 16, LIGHT_GRAY)
add_para(tf, "PB загвар (1р, 2р, 3р) аль нь явж байгааг таних", 16, LIGHT_GRAY)
add_para(tf, "Fibo 0-1-2 татаж бүтэц тодорхойлох", 16, WHITE, True)

tf2 = add_text(slide, Inches(8.5), Inches(y + 0.5), Inches(6), Inches(1.3),
               "Гол зорилго:", 16, GREEN, True)
add_para(tf2, '"Бүтэц юу харуулж байна?"', 18, WHITE, True)
add_para(tf2, "Шилжүүлэгч + PB загвар = ямар бүтэц", 15, GRAY)
add_para(tf2, "= Арилжааны БҮСИЙГ тодорхойлно", 15, GRAY)

# Arrow
add_arrow(slide, Inches(7.7), Inches(y + 2.05), Inches(0.5), Inches(0.4), GREEN)

# Layer 3 - TBM
y = 6.8
card = add_shape(slide, Inches(0.8), Inches(y), Inches(14.4), Inches(1.8), BG_CARD, ORANGE)
add_text(slide, Inches(1.2), Inches(y + 0.15), Inches(5), Inches(0.5),
         "ДАВХАР 3: TBM — ХЭЗЭЭ?", 22, ORANGE, True)
tf = add_text(slide, Inches(1.2), Inches(y + 0.65), Inches(6.5), Inches(1),
              "Өмнөх хөдөлгөөнөөс MOM олох", 16, LIGHT_GRAY)
add_para(tf, "Fibo-н 61.8 = Entry Zone", 16, WHITE, True)
add_para(tf, "Уялдаа (confluence) шалгах → ★ Entry", 16, LIGHT_GRAY)

tf2 = add_text(slide, Inches(8.5), Inches(y + 0.4), Inches(6), Inches(1.2),
               "Гол зорилго:", 16, ORANGE, True)
add_para(tf2, '"Яг хэзээ, яг хаана орох вэ?"', 18, WHITE, True)
add_para(tf2, "MOM → 61.8 → Entry = НАРИЙН ЦЭГ", 15, GRAY)

# ══════════════════════════════════════════════════════════════
# SLIDE 3: DEGI Detail
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(14), Inches(0.8),
         "ДАВХАР 1: DEGI — HTF CONTEXT", 36, CYAN, True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(5), Pt(2), CYAN)

# Left - Process
card = add_shape(slide, Inches(0.8), Inches(1.6), Inches(7), Inches(6.8), BG_CARD, CYAN)
add_text(slide, Inches(1.2), Inches(1.8), Inches(6), Inches(0.5),
         "HTF ШИНЖИЛГЭЭ (1H, 4H, Daily)", 20, CYAN, True)

steps = [
    ("1.", "Баталгаажсан candlestick", "Өмнөх лааны high/low-г эвдсэн, чиглэл баталгаажсан лаа"),
    ("2.", "Impulse таних", "Баталгаажсан candle нэг чиглэлд нэгнээс нэг чиглэлд эвдэж явна = impulse"),
    ("3.", "Pullback хүлээх", "Impulse-н эсрэг чиглэлийн хөдөлгөөн = pullback"),
    ("4.", "IDM (Inducement) sweep", "Pullback нь өмнөх лааны high/low-г хамах = liquidity хамалт"),
    ("5.", "Orderflow баталгаа", "Sweep хийсний дараа LTF дээр choch/BOS гарах"),
    ("6.", "Entry!", "LTF orderflow баталгаажсан → арилжаанд ор"),
]

y_pos = 2.5
for num, title, desc in steps:
    add_text(slide, Inches(1.3), Inches(y_pos), Inches(0.5), Inches(0.4), num, 18, GOLD, True)
    add_text(slide, Inches(1.8), Inches(y_pos), Inches(3.5), Inches(0.4), title, 17, WHITE, True)
    add_text(slide, Inches(1.8), Inches(y_pos + 0.35), Inches(5.5), Inches(0.4), desc, 13, GRAY)
    y_pos += 0.85

# Right - Entry Grades
card2 = add_shape(slide, Inches(8.2), Inches(1.6), Inches(7), Inches(3.2), BG_CARD, GOLD)
add_text(slide, Inches(8.6), Inches(1.8), Inches(6), Inches(0.5),
         "ENTRY GRADES", 22, GOLD, True)

grades = [
    ("Grade A.1", "(Хамгийн сайн)", "Sweep + HTF finalize + LTF CHOCH", GREEN),
    ("Grade A", "(Сайн)", "Sweep + HTF alignment + LTF CHOCH", BLUE),
    ("Potential A.1", "(Болзошгүй)", "Sweep + LTF CHOCH, HTF finalize байхгүй", ORANGE),
]

y_pos = 2.5
for grade, qual, desc, col in grades:
    add_text(slide, Inches(8.7), Inches(y_pos), Inches(2.5), Inches(0.35), grade, 17, col, True)
    add_text(slide, Inches(11.2), Inches(y_pos), Inches(2), Inches(0.35), qual, 14, GRAY)
    add_text(slide, Inches(8.7), Inches(y_pos + 0.3), Inches(6), Inches(0.35), desc, 13, LIGHT_GRAY)
    y_pos += 0.7

# Right bottom - Key rules
card3 = add_shape(slide, Inches(8.2), Inches(5.2), Inches(7), Inches(3.2), BG_CARD, RGBColor(0x44, 0x44, 0x66))
add_text(slide, Inches(8.6), Inches(5.4), Inches(6), Inches(0.5),
         "ЧУХАЛ ДҮРМҮҮД", 20, WHITE, True)

rules = [
    "Том цагийн баталгаа = 15min дээр Fractal шилжүүлэгч",
    "Balance бүс = жижиг цагийн orderflow",
    "Liquidity авсан бол зайлшгүй эсрэг хөдөлгөөн орно",
    "Wick хийсэн = liquidity хамалт, биеэр эвдсэн = orderflow",
    "Trend is your friend — зөвхөн трэнд дагаж арилжаал",
]

y_pos = 5.95
for rule in rules:
    add_text(slide, Inches(8.7), Inches(y_pos), Inches(6.3), Inches(0.4),
             "•  " + rule, 13, LIGHT_GRAY)
    y_pos += 0.4

# ══════════════════════════════════════════════════════════════
# SLIDE 4: FRACTAL - 3 Patterns Overview
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(14), Inches(0.8),
         "ДАВХАР 2: FRACTAL — ШИЛЖҮҮЛЭГЧИЙН 3 ЗАГВАР", 36, GREEN, True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(7), Pt(2), GREEN)

add_text(slide, Inches(0.8), Inches(1.4), Inches(14), Inches(0.5),
         "Бүс → Шилжүүлэгч → PB Загвар → Entry  |  Шилжүүлэгч бүрийн дараа PB 3 загварын аль нэгээр явна",
         16, GRAY)

# Pattern 1
card = add_shape(slide, Inches(0.5), Inches(2.1), Inches(4.8), Inches(6.3), BG_CARD, GREEN)
add_text(slide, Inches(0.8), Inches(2.3), Inches(4.2), Inches(0.5),
         "1р ЗАГВАР", 24, GREEN, True, PP_ALIGN.CENTER)
add_text(slide, Inches(0.8), Inches(2.8), Inches(4.2), Inches(0.4),
         "Ranging → Displacement → Pullback", 14, GOLD, False, PP_ALIGN.CENTER)

# Pattern 1 description
tf = add_text(slide, Inches(0.8), Inches(3.4), Inches(4.2), Inches(4.5),
              "Бүтэц:", 16, WHITE, True)
add_para(tf, "1. Бүс (ranging zone) үүснэ", 14, LIGHT_GRAY)
add_para(tf, "2. Том лаа бүсийг эвдэнэ (шилжүүлэгч)", 14, LIGHT_GRAY)
add_para(tf, "3. Шууд pullback орно", 14, LIGHT_GRAY)
add_para(tf, "4. Pullback дотор ranging бүс үүснэ", 14, LIGHT_GRAY)
add_para(tf, "5. Дотоод бүсээс эсрэг displacement", 14, LIGHT_GRAY)
add_para(tf, "6. Entry!", 14, GREEN, True)
add_para(tf, "", 10, GRAY)
add_para(tf, "Онцлог:", 15, GOLD, True)
add_para(tf, "PB дотроо бүс → шилжүүлэгч давтана", 13, GRAY)
add_para(tf, "= Fractal-ийн Fractal", 13, GRAY)
add_para(tf, "", 10, GRAY)
add_para(tf, "candle-butesh2.png дээрх зураг", 12, RGBColor(0x88, 0x88, 0x88))

# Pattern 2
card = add_shape(slide, Inches(5.6), Inches(2.1), Inches(4.8), Inches(6.3), BG_CARD, BLUE)
add_text(slide, Inches(5.9), Inches(2.3), Inches(4.2), Inches(0.5),
         "2р ЗАГВАР", 24, BLUE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(5.9), Inches(2.8), Inches(4.2), Inches(0.4),
         "Simple Pullback (1-3 давалгаа)", 14, GOLD, False, PP_ALIGN.CENTER)

tf = add_text(slide, Inches(5.9), Inches(3.4), Inches(4.2), Inches(4.5),
              "Бүтэц:", 16, WHITE, True)
add_para(tf, "1. Шилжүүлэгч гарна", 14, LIGHT_GRAY)
add_para(tf, "2. Энгийн pullback (1-3 давалгаа)", 14, LIGHT_GRAY)
add_para(tf, "3. Дотроо ranging/displacement ҮГҮЙ", 14, LIGHT_GRAY)
add_para(tf, "4. Шууд эргээд entry чиглэлд явна", 14, LIGHT_GRAY)
add_para(tf, "5. Entry!", 14, BLUE, True)
add_para(tf, "", 10, GRAY)
add_para(tf, "Онцлог:", 15, GOLD, True)
add_para(tf, "Хамгийн энгийн, хурдан загвар", 13, GRAY)
add_para(tf, "Бүс доторх displacement байхгүй", 13, GRAY)
add_para(tf, "Шууд V эсвэл W хэлбэрийн PB", 13, GRAY)
add_para(tf, "", 10, GRAY)
add_para(tf, "candle-butesh.png дээрх зураг", 12, RGBColor(0x88, 0x88, 0x88))

# Pattern 3
card = add_shape(slide, Inches(10.7), Inches(2.1), Inches(4.8), Inches(6.3), BG_CARD, PURPLE)
add_text(slide, Inches(11), Inches(2.3), Inches(4.2), Inches(0.5),
         "3р ЗАГВАР", 24, PURPLE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(11), Inches(2.8), Inches(4.2), Inches(0.4),
         "Complex Pullback (4+ давалгаа)", 14, GOLD, False, PP_ALIGN.CENTER)

tf = add_text(slide, Inches(11), Inches(3.4), Inches(4.2), Inches(4.5),
              "Бүтэц:", 16, WHITE, True)
add_para(tf, "1. Шилжүүлэгч гарна", 14, LIGHT_GRAY)
add_para(tf, "2. Олон давалгаатай pullback (4+)", 14, LIGHT_GRAY)
add_para(tf, "3. Масштаб өөрчлөгдөж явна", 14, LIGHT_GRAY)
add_para(tf, "4. Том → жижиг → том давалгаа", 14, LIGHT_GRAY)
add_para(tf, "5. Entry!", 14, PURPLE, True)
add_para(tf, "", 10, GRAY)
add_para(tf, "Онцлог:", 15, GOLD, True)
add_para(tf, "Хамгийн удаан, олон давалгаатай", 13, GRAY)
add_para(tf, "Scale-ээр Fractal давтагдана", 13, GRAY)
add_para(tf, "Тэвчээр шаардлагатай", 13, GRAY)
add_para(tf, "", 10, GRAY)
add_para(tf, "candle-butesh3.png дээрх зураг", 12, RGBColor(0x88, 0x88, 0x88))

# ══════════════════════════════════════════════════════════════
# SLIDE 5: PB Structure Deep Dive
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(14), Inches(0.8),
         "PB БҮТЭЦ — НАРИЙН ЗАДАРГАА", 36, GREEN, True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(5), Pt(2), GREEN)

# LEFT: PB anatomy
card = add_shape(slide, Inches(0.5), Inches(1.5), Inches(7.3), Inches(7), BG_CARD, GREEN)
add_text(slide, Inches(0.9), Inches(1.7), Inches(6.5), Inches(0.5),
         "PULLBACK-ИЙН АНАТОМИ", 22, GREEN, True)

tf = add_text(slide, Inches(0.9), Inches(2.3), Inches(6.5), Inches(6),
              "Шилжүүлэгч → PB эхлэл", 17, GOLD, True)
add_para(tf, "Шилжүүлэгч = том лаа бүсийг биеэрээ эвдэх", 14, LIGHT_GRAY)
add_para(tf, "Биеэрээ хаагдсан байх ЁСТОЙ (wick биш)", 14, WHITE, True)
add_para(tf, "Хэмжээсийн урт (0-1) богино = шилжүүлэгч биш", 14, LIGHT_GRAY)
add_para(tf, "", 8, GRAY)

add_para(tf, "PB эхлэл таних:", 17, GOLD, True)
add_para(tf, "Шилжүүлэгчийн дараа эсрэг чиглэлийн лаа = PB эхэлсэн", 14, LIGHT_GRAY)
add_para(tf, "PB эхлэхгүй бол 10 bar-н дараа шилжүүлэгч хүчингүй", 14, RED)
add_para(tf, "", 8, GRAY)

add_para(tf, "PB дотоод бүтэц:", 17, GOLD, True)
add_para(tf, "Давалгаа тоолох — эсрэг чиглэл солих бүрд +1", 14, LIGHT_GRAY)
add_para(tf, "1-3 давалгаа = 2р загвар (энгийн)", 14, BLUE)
add_para(tf, "4+ давалгаа = 3р загвар (complex)", 14, PURPLE)
add_para(tf, "Дотроо ranging + displacement = 1р загвар", 14, GREEN)
add_para(tf, "", 8, GRAY)

add_para(tf, "PB доторх RANGING шалгах:", 17, GOLD, True)
add_para(tf, "5 bar-н range < ATR x 2.5 → ranging эхэлсэн", 14, LIGHT_GRAY)
add_para(tf, "Ranging + эсрэг displacement = 1р загвар баталгаа", 14, LIGHT_GRAY)

# RIGHT: PB end detection
card = add_shape(slide, Inches(8.1), Inches(1.5), Inches(7.3), Inches(7), BG_CARD, ORANGE)
add_text(slide, Inches(8.5), Inches(1.7), Inches(6.5), Inches(0.5),
         "PB ТӨГСГӨЛ ОЛОХ = TBM", 22, ORANGE, True)

tf = add_text(slide, Inches(8.5), Inches(2.3), Inches(6.5), Inches(6),
              "MOM олох (2 боломж):", 17, GOLD, True)
add_para(tf, "", 6, GRAY)
add_para(tf, "Боломж 1: Үндсэн Low/High = MOM", 16, WHITE, True)
add_para(tf, "Өмнөх impulse-н 61.8-г сэтлэхэд тухайн", 14, LIGHT_GRAY)
add_para(tf, "impulse-н эхлэл low (bull) / high (bear) = MOM", 14, LIGHT_GRAY)
add_para(tf, "", 6, GRAY)

add_para(tf, "Боломж 2: Жижиг impulse-н 61.8 = MOM", 16, WHITE, True)
add_para(tf, "Сэтлэсэн хөдөлгөөн дотоод жижиг impulse", 14, LIGHT_GRAY)
add_para(tf, "үүсгэсэн → тэр жижиг imp-н 61.8 = MOM", 14, LIGHT_GRAY)
add_para(tf, "", 8, GRAY)

add_para(tf, "Энэ 2-н АЛЬ НЭГ нь MOM болно!", 16, GOLD, True)
add_para(tf, "", 8, GRAY)

add_para(tf, "61.8 Entry Zone:", 17, GOLD, True)
add_para(tf, "MOM олсны дараа Fibo 0-1 татна", 14, LIGHT_GRAY)
add_para(tf, "61.8 = PB-н яг төгсөх цэг = ENTRY ZONE", 14, GREEN, True)
add_para(tf, "61.8 - 72% хооронд = хамгийн сайн бүс", 14, LIGHT_GRAY)
add_para(tf, "", 8, GRAY)

add_para(tf, "Шүүлтүүрүүд:", 17, RED, True)
add_para(tf, "Эрлийз: MOM ↔ 61.8 зай > ATR x 5 → ХАЯХ", 14, RED)
add_para(tf, "Уялдаа: Олон MOM-н 61.8 давхцах → ★ ENTRY", 14, GREEN)

# ══════════════════════════════════════════════════════════════
# SLIDE 6: PB Wave Structure Detail
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(14), Inches(0.8),
         "PB ДАВАЛГААНЫ БҮТЭЦ — НАРИЙН", 36, GREEN, True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(5), Pt(2), GREEN)

# Bull PB example
card = add_shape(slide, Inches(0.5), Inches(1.5), Inches(7.3), Inches(3.3), BG_CARD, GREEN)
add_text(slide, Inches(0.9), Inches(1.7), Inches(6), Inches(0.5),
         "BULL SETUP — PB БҮТЭЦ", 20, GREEN, True)

tf = add_text(slide, Inches(0.9), Inches(2.3), Inches(6.5), Inches(2.3),
              "Хөдөлгөөний дараалал (дээшээ):", 15, GOLD, True)
add_para(tf, "", 4, GRAY)
add_para(tf, "[ БҮС ]  →  ▲ ШИЛЖҮҮЛЭГЧ (том ногоон лаа, бүсийг дээш эвднэ)", 13, GREEN)
add_para(tf, "           →  ▼ PB эхлэл (улаан лаа гарна)", 13, RED)
add_para(tf, "           →  ▼▲▼ Давалгаа (1р, 2р, 3р загварын аль нэг)", 13, LIGHT_GRAY)
add_para(tf, "           →  IDM sweep (өмнөх low-г хамна = Degi)", 13, CYAN)
add_para(tf, "           →  MOM олдоно (TBM)", 13, ORANGE)
add_para(tf, "           →  61.8 дээр ногоон лаа = ▲ ENTRY", 13, GREEN, True)

# Bear PB example
card = add_shape(slide, Inches(0.5), Inches(5.1), Inches(7.3), Inches(3.3), BG_CARD, RED)
add_text(slide, Inches(0.9), Inches(5.3), Inches(6), Inches(0.5),
         "BEAR SETUP — PB БҮТЭЦ", 20, RED, True)

tf = add_text(slide, Inches(0.9), Inches(5.9), Inches(6.5), Inches(2.3),
              "Хөдөлгөөний дараалал (доошоо):", 15, GOLD, True)
add_para(tf, "", 4, GRAY)
add_para(tf, "[ БҮС ]  →  ▼ ШИЛЖҮҮЛЭГЧ (том улаан лаа, бүсийг доош эвднэ)", 13, RED)
add_para(tf, "           →  ▲ PB эхлэл (ногоон лаа гарна)", 13, GREEN)
add_para(tf, "           →  ▲▼▲ Давалгаа (1р, 2р, 3р загварын аль нэг)", 13, LIGHT_GRAY)
add_para(tf, "           →  IDM sweep (өмнөх high-г хамна = Degi)", 13, CYAN)
add_para(tf, "           →  MOM олдоно (TBM)", 13, ORANGE)
add_para(tf, "           →  61.8 дээр улаан лаа = ▼ ENTRY", 13, RED, True)

# Right side - Validation rules
card = add_shape(slide, Inches(8.1), Inches(1.5), Inches(7.3), Inches(3.3), BG_CARD, GOLD)
add_text(slide, Inches(8.5), Inches(1.7), Inches(6.5), Inches(0.5),
         "ШИЛЖҮҮЛЭГЧ ШАЛГАХ", 20, GOLD, True)

tf = add_text(slide, Inches(8.5), Inches(2.3), Inches(6.5), Inches(2.3),
              "", 14, WHITE)
rules = [
    ("✓", "Биеэрээ хаагдсан байх (wick-ээр биш)", GREEN),
    ("✓", "Хэмжээс (0-1) хангалттай урт байх", GREEN),
    ("✓", "Өмнөх бүсийн range-с том displacement", GREEN),
    ("✓", "Огцом хөдөлгөөн (consolidation-гүй)", GREEN),
    ("✗", "Удаан хугацаанд огцомгүй = шилжүүлэгч биш", RED),
    ("✗", "1н цэг шалгаж, pullback хэт гүн = буруу таних", RED),
    ("✗", "Өмнөх жижиг pullback-с хүчилж татах = биш", RED),
]
for mark, rule, col in rules:
    p = tf.add_paragraph()
    p.text = f"  {mark}  {rule}"
    p.font.size = Pt(13)
    p.font.color.rgb = col
    p.space_before = Pt(4)

# Right bottom - MTF
card = add_shape(slide, Inches(8.1), Inches(5.1), Inches(7.3), Inches(3.3), BG_CARD, CYAN)
add_text(slide, Inches(8.5), Inches(5.3), Inches(6.5), Inches(0.5),
         "ЦАГ ХООРОНДЫН ХОЛБОО", 20, CYAN, True)

tf = add_text(slide, Inches(8.5), Inches(5.9), Inches(6.5), Inches(2.3),
              "Scalp:  1-5min entry → 30min, 1H харах", 15, LIGHT_GRAY)
add_para(tf, "Day:    5-15min entry → 2H, 4H харах", 15, LIGHT_GRAY)
add_para(tf, "Swing:  15-30min entry → 5H, Daily харах", 15, LIGHT_GRAY)
add_para(tf, "", 8, GRAY)
add_para(tf, "Том цагийн 2 лаа = Жижиг цагийн шилжүүлэгч", 15, WHITE, True)
add_para(tf, "1H дээр 2 ногоон лаа = 15min дээр displacement", 14, LIGHT_GRAY)
add_para(tf, "→ Энэ бол Degi ↔ Fractal-ийн ХОЛБООС", 14, CYAN, True)

# ══════════════════════════════════════════════════════════════
# SLIDE 7: TBM - MOM + 61.8 Detail
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(14), Inches(0.8),
         "ДАВХАР 3: TBM — MOM + FIBO 61.8", 36, ORANGE, True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(5), Pt(2), ORANGE)

# MOM Finding Process
card = add_shape(slide, Inches(0.5), Inches(1.5), Inches(7.3), Inches(7), BG_CARD, ORANGE)
add_text(slide, Inches(0.9), Inches(1.7), Inches(6.5), Inches(0.5),
         "MOM ОЛОХ ПРОЦЕСС", 22, ORANGE, True)

tf = add_text(slide, Inches(0.9), Inches(2.3), Inches(6.5), Inches(5.8),
              "Алхам 1: Impulse таних", 16, GOLD, True)
add_para(tf, "Swing High → Swing Low (эсвэл эсрэгээр)", 14, LIGHT_GRAY)
add_para(tf, "Range > ATR x 2 байх ёстой (жижиг импулс биш)", 14, LIGHT_GRAY)
add_para(tf, "", 6, GRAY)

add_para(tf, "Алхам 2: 61.8 сэтлэгдэх", 16, GOLD, True)
add_para(tf, "Тухайн impulse-н 61.8-г ханш сэтлэх = трэнд эргэлт", 14, LIGHT_GRAY)
add_para(tf, "Биеэр сэтлэх ёстой (wick биш)", 14, WHITE, True)
add_para(tf, "Bearish candle (close < open) байх ёстой", 14, LIGHT_GRAY)
add_para(tf, "", 6, GRAY)

add_para(tf, "Алхам 3: MOM тодорхойлох", 16, GOLD, True)
add_para(tf, "БОЛОМЖ 1:", 14, GREEN, True)
add_para(tf, "  Тухайн impulse-н High (bear) / Low (bull) = MOM", 13, LIGHT_GRAY)
add_para(tf, "  → Энэ бол үндсэн low, эсрэг хөдөлгөөн эхлэх цэг", 13, GRAY)
add_para(tf, "", 4, GRAY)
add_para(tf, "БОЛОМЖ 2:", 14, BLUE, True)
add_para(tf, "  Жижиг impulse-н 61.8 = MOM", 13, LIGHT_GRAY)
add_para(tf, "  → Сэтлэсэн хөдөлгөөн доторх жижиг imp-н 61.8", 13, GRAY)
add_para(tf, "", 6, GRAY)

add_para(tf, "Алхам 4: Fibo татах", 16, GOLD, True)
add_para(tf, "0 = Impulse-н нэг үзүүр", 14, LIGHT_GRAY)
add_para(tf, "1 = Impulse-н нөгөө үзүүр", 14, LIGHT_GRAY)
add_para(tf, "MOM = Momentum цэг (улаан шугам)", 14, RED)
add_para(tf, "61.8 = ENTRY ZONE (улаан зузаан шугам)", 14, RED, True)

# Right - Entry Logic
card = add_shape(slide, Inches(8.1), Inches(1.5), Inches(7.3), Inches(3.5), BG_CARD, GREEN)
add_text(slide, Inches(8.5), Inches(1.7), Inches(6.5), Inches(0.5),
         "61.8 ENTRY LOGIC", 22, GREEN, True)

tf = add_text(slide, Inches(8.5), Inches(2.3), Inches(6.5), Inches(2.5),
              "Bull Entry (bear fibo-н 61.8 дээр):", 16, GREEN, True)
add_para(tf, "Low ≤ 61.8  (61.8-д хүрсэн)", 14, LIGHT_GRAY)
add_para(tf, "Close > 61.8  (дээш буцсан)", 14, LIGHT_GRAY)
add_para(tf, "Close > Open  (ногоон лаа)", 14, LIGHT_GRAY)
add_para(tf, "= ▲ ENTRY", 16, GREEN, True)
add_para(tf, "", 6, GRAY)
add_para(tf, "Bear Entry (bull fibo-н 61.8 дээр):", 16, RED, True)
add_para(tf, "High ≥ 61.8  →  Close < 61.8  →  Close < Open", 14, LIGHT_GRAY)
add_para(tf, "= ▼ ENTRY", 16, RED, True)

# Right - Confluence
card = add_shape(slide, Inches(8.1), Inches(5.3), Inches(7.3), Inches(3.2), BG_CARD, GOLD)
add_text(slide, Inches(8.5), Inches(5.5), Inches(6.5), Inches(0.5),
         "УЯЛДАА + ЭРЛИЙЗ", 22, GOLD, True)

tf = add_text(slide, Inches(8.5), Inches(6.1), Inches(6.5), Inches(2.2),
              "★ Уялдаа (Confluence):", 16, GREEN, True)
add_para(tf, "Олон MOM-ийн 61.8-72% бүс давхцах", 14, LIGHT_GRAY)
add_para(tf, "= Маш хүчтэй entry zone", 14, LIGHT_GRAY)
add_para(tf, "= ★ ENTRY (тодоор харагдана)", 14, GREEN, True)
add_para(tf, "", 6, GRAY)

add_para(tf, "✗ Эрлийз (Reject):", 16, RED, True)
add_para(tf, "MOM ↔ 61.8 хоорондох зай > ATR x 5", 14, LIGHT_GRAY)
add_para(tf, "= Fibo хэт сунасан, итгэх боломжгүй → ХАЯХ", 14, RED)

# ══════════════════════════════════════════════════════════════
# SLIDE 8: Full Flow Diagram
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(14), Inches(0.8),
         "БҮТЭН ПРОЦЕСС — A.1 ENTRY SETUP", 36, GOLD, True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(5), Pt(2), GOLD)

# Step boxes flowing left to right
steps_data = [
    ("1", "HTF ЧИГЛЭЛ", "4H/1H дээр\nтрэнд таних", CYAN, BG_CARD),
    ("2", "LIQUIDITY", "IDM sweep\nхүлээх", CYAN, BG_CARD),
    ("3", "БҮС ОЛОХ", "Ranging zone\nтодорхойлох", GREEN, BG_CARD),
    ("4", "ШИЛЖҮҮЛЭГЧ", "Том лаа\nбүсийг эвдэх", GREEN, BG_CARD),
    ("5", "PB ЗАГВАР", "1р/2р/3р\nзагвар таних", GREEN, BG_CARD),
    ("6", "MOM ОЛОХ", "Өмнөх imp-с\nMOM олох", ORANGE, BG_CARD),
    ("7", "61.8 ENTRY", "Entry zone\nхүлээх", ORANGE, BG_CARD),
    ("8", "★ ENTRY", "Уялдаа +\nCandle = ОР!", GOLD, BG_CARD),
]

for i, (num, title, desc, col, bg) in enumerate(steps_data):
    x = Inches(0.3 + i * 1.95)
    y = Inches(1.8)

    card = add_shape(slide, x, y, Inches(1.75), Inches(2.2), bg, col)
    add_text(slide, x, Inches(1.85), Inches(1.75), Inches(0.4), num, 28, col, True, PP_ALIGN.CENTER)
    add_text(slide, x, Inches(2.25), Inches(1.75), Inches(0.4), title, 12, col, True, PP_ALIGN.CENTER)

    # desc lines
    lines = desc.split('\n')
    for j, line in enumerate(lines):
        add_text(slide, x, Inches(2.7 + j * 0.3), Inches(1.75), Inches(0.3),
                 line, 11, GRAY, False, PP_ALIGN.CENTER)

    # Arrow between steps
    if i < len(steps_data) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                        Emu(int((0.3 + (i+1) * 1.95) * 914400 - Inches(0.15).emu)),
                                        Inches(2.7),
                                        Inches(0.25), Inches(0.2))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0x44, 0x44, 0x66)
        arrow.line.fill.background()

# Bottom section - Checklist
card = add_shape(slide, Inches(0.5), Inches(4.5), Inches(15), Inches(4), BG_CARD, GOLD)
add_text(slide, Inches(0.9), Inches(4.7), Inches(14), Inches(0.5),
         "ENTRY CHECKLIST — Арилжаанд орохын өмнөх шалгах жагсаалт", 22, GOLD, True)

# Column 1
tf = add_text(slide, Inches(0.9), Inches(5.3), Inches(4.5), Inches(3),
              "DEGI (HTF Context)", 16, CYAN, True)
add_para(tf, "☐ HTF трэнд тодорхой юу?", 14, LIGHT_GRAY)
add_para(tf, "☐ Liquidity sweep болсон уу?", 14, LIGHT_GRAY)
add_para(tf, "☐ IDM хамагдсан уу?", 14, LIGHT_GRAY)
add_para(tf, "☐ HTF choch/BOS баталгаажсан уу?", 14, LIGHT_GRAY)
add_para(tf, "☐ Entry Grade (A.1 / A / Pot.A.1)?", 14, LIGHT_GRAY)

# Column 2
tf = add_text(slide, Inches(5.7), Inches(5.3), Inches(4.5), Inches(3),
              "FRACTAL (Бүтэц)", 16, GREEN, True)
add_para(tf, "☐ Бүс тодорхойлогдсон уу?", 14, LIGHT_GRAY)
add_para(tf, "☐ Шилжүүлэгч биеэр хаагдсан уу?", 14, LIGHT_GRAY)
add_para(tf, "☐ PB загвар аль нь вэ? (1р/2р/3р)", 14, LIGHT_GRAY)
add_para(tf, "☐ Fibo 0-1-2 зөв татсан уу?", 14, LIGHT_GRAY)
add_para(tf, "☐ Цаг шилжүүлэл зөв үү?", 14, LIGHT_GRAY)

# Column 3
tf = add_text(slide, Inches(10.5), Inches(5.3), Inches(4.5), Inches(3),
              "TBM (Entry)", 16, ORANGE, True)
add_para(tf, "☐ MOM олдсон уу? (Боломж 1/2)", 14, LIGHT_GRAY)
add_para(tf, "☐ Эрлийз биш үү? (MOM↔61.8 зай)", 14, LIGHT_GRAY)
add_para(tf, "☐ 61.8 zone дээр ирсэн үү?", 14, LIGHT_GRAY)
add_para(tf, "☐ Уялдаа байна уу? (★)", 14, LIGHT_GRAY)
add_para(tf, "☐ Entry candle баталгаажсан уу?", 14, LIGHT_GRAY)

# ══════════════════════════════════════════════════════════════
# SLIDE 9: Candlestick Patterns for Entry
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(14), Inches(0.8),
         "ЛААНЫ БҮТЭЦ — ENTRY БАТАЛГАА", 36, GOLD, True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(5), Pt(2), GOLD)

# Candlestick shift
card = add_shape(slide, Inches(0.5), Inches(1.5), Inches(4.8), Inches(3.2), BG_CARD, GREEN)
add_text(slide, Inches(0.9), Inches(1.7), Inches(4), Inches(0.5),
         "CANDLESTICK SHIFT", 20, GREEN, True)
tf = add_text(slide, Inches(0.9), Inches(2.3), Inches(4), Inches(2.2),
              "2 лааны шилжилт хэсгийг ажигла:", 14, GOLD, True)
add_para(tf, "Төгсгөл → Дараагийн лааны эхлэл", 14, LIGHT_GRAY)
add_para(tf, "Унаж орж ирвэл → дараагийх ногоон", 14, GREEN)
add_para(tf, "Өсөлт орж ирвэл → дараагийх улаан", 14, RED)
add_para(tf, "", 6, GRAY)
add_para(tf, "Шилжилтийг Fibo дээр ч харна", 14, LIGHT_GRAY)

# Лааны Өр
card = add_shape(slide, Inches(5.6), Inches(1.5), Inches(4.8), Inches(3.2), BG_CARD, BLUE)
add_text(slide, Inches(6), Inches(1.7), Inches(4), Inches(0.5),
         "ЛААНЫ ӨР", 20, BLUE, True)
tf = add_text(slide, Inches(6), Inches(2.3), Inches(4), Inches(2.2),
              "Ханш тэгш хэмжээгээр хөдөлдөг:", 14, GOLD, True)
add_para(tf, "Өмнөх хэмжээгээр дахин давтана", 14, LIGHT_GRAY)
add_para(tf, "Лааны өндрийг аль нэг тал руу давтах", 14, LIGHT_GRAY)
add_para(tf, "4H+ дээр харах → шууд бодоод байхгүй", 14, LIGHT_GRAY)
add_para(tf, "Fibo дээр давхцаад гараад ирдэг", 14, LIGHT_GRAY)

# Inside Bar
card = add_shape(slide, Inches(10.7), Inches(1.5), Inches(4.8), Inches(3.2), BG_CARD, PURPLE)
add_text(slide, Inches(11.1), Inches(1.7), Inches(4), Inches(0.5),
         "INSIDE BAR (IB)", 20, PURPLE, True)
tf = add_text(slide, Inches(11.1), Inches(2.3), Inches(4), Inches(2.2),
              "Голын candle-н дотор багтсан:", 14, GOLD, True)
add_para(tf, "2 хажуу талын H/L эвдэж чадахгүй", 14, LIGHT_GRAY)
add_para(tf, "3дахь лаа эвдэх тал руу явна", 14, LIGHT_GRAY)
add_para(tf, "3дахь нь 1ээсээ заавал том байх", 14, WHITE, True)
add_para(tf, "IB + Fibo = маш сайн ашиглагдана", 14, LIGHT_GRAY)

# Bottom - Key Level + Liquidity
card = add_shape(slide, Inches(0.5), Inches(5), Inches(7.3), Inches(3.5), BG_CARD, CYAN)
add_text(slide, Inches(0.9), Inches(5.2), Inches(6.5), Inches(0.5),
         "KEY LEVEL + LIQUIDITY", 20, CYAN, True)
tf = add_text(slide, Inches(0.9), Inches(5.8), Inches(6.5), Inches(2.5),
              "Key Level:", 16, WHITE, True)
add_para(tf, "Ханш ойлт авсан бүс = дэмжлэг/эсэргүүцэл", 14, LIGHT_GRAY)
add_para(tf, "61.8 zone = маш сайн key level", 14, LIGHT_GRAY)
add_para(tf, "3+ удаа ойсон бүс = хүчтэй", 14, LIGHT_GRAY)
add_para(tf, "", 6, GRAY)
add_para(tf, "Liquidity sweep:", 16, WHITE, True)
add_para(tf, "Өмнөх лааны low эвдсэн = шууд эсрэг хөдөлгөөн", 14, LIGHT_GRAY)
add_para(tf, "Wick = liquidity хамалт, Биеэр = orderflow", 14, GOLD, True)

# Bottom right - Сэгрэлт
card = add_shape(slide, Inches(8.1), Inches(5), Inches(7.3), Inches(3.5), BG_CARD, RED)
add_text(slide, Inches(8.5), Inches(5.2), Inches(6.5), Inches(0.5),
         "СЭГРЭЛТ + АНХААРУУЛГА", 20, RED, True)
tf = add_text(slide, Inches(8.5), Inches(5.8), Inches(6.5), Inches(2.5),
              "Сэгрэлт:", 16, WHITE, True)
add_para(tf, "Орой/ёроол олон удаа эвдрэх = сэгрэлт", 14, LIGHT_GRAY)
add_para(tf, "Сэтлээд явчхаараа гэж бодох хэрэггүй", 14, LIGHT_GRAY)
add_para(tf, "Арилжаа НЭЭХГҮЙ", 14, RED, True)
add_para(tf, "", 6, GRAY)
add_para(tf, "Range үед:", 16, WHITE, True)
add_para(tf, "Range задарсны дараа л арилжаа хий", 14, LIGHT_GRAY)
add_para(tf, "Range дотор трэнд эргэлт fibo бүү хэрэглэ", 14, RED)

# ══════════════════════════════════════════════════════════════
# SLIDE 10: Summary / Recap
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(14), Inches(0.8),
         "ДҮГНЭЛТ", 42, GOLD, True, PP_ALIGN.CENTER)
add_shape(slide, Inches(5), Inches(1.15), Inches(6), Pt(2), GOLD)

# Core principle
card = add_shape(slide, Inches(2), Inches(1.8), Inches(12), Inches(2), BG_CARD, GOLD)
tf = add_text(slide, Inches(2.5), Inches(2.0), Inches(11), Inches(1.6),
              "НЭГ ПРОЦЕССЫН 3 ZOOM LEVEL", 28, GOLD, True, PP_ALIGN.CENTER)
add_para(tf, "", 6, GRAY)
add_para(tf, "DEGI = Яагаад энэ чиглэлд?  (HTF context, liquidity)", 18, CYAN, False)
add_para(tf, "FRACTAL = Хаана арилжаалах?  (бүс, шилжүүлэгч, PB бүтэц)", 18, GREEN, False)
add_para(tf, "TBM = Хэзээ яг орох?  (MOM → 61.8 → Entry)", 18, ORANGE, False)

# Key rules
card = add_shape(slide, Inches(0.8), Inches(4.2), Inches(7), Inches(4.3), BG_CARD, WHITE)
add_text(slide, Inches(1.2), Inches(4.4), Inches(6), Inches(0.5),
         "АЛТАН ДҮРМҮҮД", 22, GOLD, True)
tf = add_text(slide, Inches(1.2), Inches(5), Inches(6.3), Inches(3.3),
              "1.  Том цагийн баталгаагүй бол бага цаг луугаа ор", 15, WHITE, True)
add_para(tf, "2.  Шилжүүлэгч биеэрээ хаагдсан байх ёстой", 15, WHITE, True)
add_para(tf, "3.  61.8 хүрээгүй бол жинхэнэ pullback биш", 15, WHITE, True)
add_para(tf, "4.  Эрлийз fibo-г хаях (MOM↔61.8 хэт хол)", 15, WHITE, True)
add_para(tf, "5.  Уялдаа = хамгийн итгэлтэй entry (★)", 15, WHITE, True)
add_para(tf, "6.  Range үед битгий арилжаалаарай", 15, WHITE, True)
add_para(tf, "7.  Trend is your friend — зөвхөн трэнд дагаж ор", 15, WHITE, True)
add_para(tf, "8.  Wick = liquidity, Биеэр = orderflow", 15, WHITE, True)

# Remember box
card = add_shape(slide, Inches(8.2), Inches(4.2), Inches(7), Inches(4.3), BG_CARD, GREEN)
add_text(slide, Inches(8.6), Inches(4.4), Inches(6), Inches(0.5),
         "САНАХ", 22, GREEN, True)
tf = add_text(slide, Inches(8.6), Inches(5), Inches(6.3), Inches(3.3),
              "3 давхар шүүлтүүр давхцвал =", 16, LIGHT_GRAY)
add_para(tf, "ХАМГИЙН ӨНДӨР МАГАДЛАЛТАЙ ENTRY", 20, GOLD, True)
add_para(tf, "", 10, GRAY)
add_para(tf, "Setup:", 16, WHITE, True)
add_para(tf, "1. Liquidity харна", 16, CYAN)
add_para(tf, "2. Fibo зөв зурагдсан байх", 16, GREEN)
add_para(tf, "3. Цаг шилжилтээ харах", 16, BLUE)
add_para(tf, "4. Лааны өр, IB, engulfing хослуулж баталж болно", 16, PURPLE)
add_para(tf, "", 10, GRAY)
add_para(tf, "Бүгд давхцвал → ENTRY!", 20, GREEN, True)

# Save
output_path = "/Users/battugs/Desktop/Project/forex-bots/Combined_Entry_System.pptx"
prs.save(output_path)
print(f"PPT saved: {output_path}")
